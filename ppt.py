from pptx import Presentation
from wand.image import Image
import os
from pptx.util import Inches
from glob import glob

#To create images with watermark
def add_watermark():
    sourceExt = 'jpg'
    files = os.listdir('Sample/')
    
    for file in files:
        filename, file_extension = os.path.splitext('Sample/' + file)
        updatedFilename = filename.replace('Sample/', '')
        if file_extension == ('.' + sourceExt):
            with Image(filename=filename + file_extension) as img:
                img.composite_channel("all_channels", Image(filename='nike_black.png'), "dissolve")
                img.save(filename= 'Watermarked/' + updatedFilename + '.' + sourceExt)
    
    
#To create ppt with images with watermark
def generate_ppt():
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[1]
    i = 1
    images = glob('Watermarked/*.JPG')

    for img in images:
        
        slide = prs.slides.add_slide(blank_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = 'Sample Title ' + str(i)

        tf = body_shape.text_frame
        tf.text = 'Sample Subtitle ' + str(i)
        top = Inches(2.5)
        left = Inches(1.8)
        height = Inches(4.5)
        pic = slide.shapes.add_picture(img, left, top, height=height)
        i += 1

    prs.save("WatermarkedPPT.pptx")
    os.startfile("WatermarkedPPT.pptx")


add_watermark()
generate_ppt()
